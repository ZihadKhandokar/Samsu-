from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.util import Inches

from files.base_adapter import FileAdapter


class PowerPointAdapter(FileAdapter):
    def create(self, path: Path, payload: Any) -> dict[str, Any]:
        if not isinstance(payload, dict):
            raise TypeError("PowerPoint creation requires a structured object.")
        presentation = Presentation()
        for slide_spec in payload.get("slides", []):
            self._add_slide(presentation, slide_spec)
        presentation.save(path)
        return {"format": "pptx", "slides": len(presentation.slides)}

    @staticmethod
    def _add_slide(presentation: Presentation, spec: dict[str, Any]) -> None:
        layout_index = int(spec.get("layout", 1))
        layout_index = min(max(layout_index, 0), len(presentation.slide_layouts) - 1)
        slide = presentation.slides.add_slide(presentation.slide_layouts[layout_index])
        if slide.shapes.title is not None:
            slide.shapes.title.text = str(spec.get("title", ""))
        body = spec.get("body")
        if body is not None:
            body_placeholder = next(
                (shape for shape in slide.placeholders if shape != slide.shapes.title and hasattr(shape, "text_frame")),
                None,
            )
            if body_placeholder is not None:
                text_frame = body_placeholder.text_frame
                text_frame.clear()
                items = body if isinstance(body, list) else [body]
                for index, item in enumerate(items):
                    paragraph = text_frame.paragraphs[0] if index == 0 else text_frame.add_paragraph()
                    paragraph.text = str(item)
        for image in spec.get("images", []):
            slide.shapes.add_picture(
                str(image["path"]),
                Inches(float(image.get("left", 1))),
                Inches(float(image.get("top", 1))),
                width=Inches(float(image.get("width", 5))),
            )

    def read(self, path: Path) -> dict[str, Any]:
        presentation = Presentation(path)
        slides = []
        for index, slide in enumerate(presentation.slides, start=1):
            texts = [shape.text for shape in slide.shapes if hasattr(shape, "text") and shape.text]
            slides.append({"number": index, "texts": texts})
        return {"format": "pptx", "slides": slides}

    def edit(self, path: Path, operations: list[dict[str, Any]]) -> dict[str, Any]:
        presentation = Presentation(path)
        for operation in operations:
            name = operation.get("operation")
            if name == "add_slide":
                self._add_slide(presentation, operation)
            elif name == "replace_text":
                old = str(operation.get("old_text", ""))
                new = str(operation.get("new_text", ""))
                replacements = 0
                for slide in presentation.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text_frame"):
                            for paragraph in shape.text_frame.paragraphs:
                                for run in paragraph.runs:
                                    if old in run.text:
                                        run.text = run.text.replace(old, new)
                                        replacements += 1
                if replacements == 0:
                    raise ValueError("PowerPoint text to replace was not found.")
            elif name == "set_title":
                slide_number = int(operation["slide"])
                slide = presentation.slides[slide_number - 1]
                if slide.shapes.title is None:
                    raise ValueError("Selected slide has no title shape.")
                slide.shapes.title.text = str(operation.get("text", ""))
            else:
                raise ValueError(f"Unsupported PowerPoint operation: {name}")
        presentation.save(path)
        return {"format": "pptx", "slides": len(presentation.slides)}

    def verify(self, path: Path) -> dict[str, Any]:
        result = super().verify(path)
        result["slides"] = len(Presentation(path).slides)
        return result
